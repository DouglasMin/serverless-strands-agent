import base64
import io
import json
import re
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from strands import tool


def _get_theme_colors(theme: str) -> dict[str, RGBColor]:
    if theme == "light":
        return {
            "bg": RGBColor(0xF8, 0xFA, 0xFC),
            "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
            "title": RGBColor(0x0F, 0x17, 0x2A),
            "body": RGBColor(0x33, 0x41, 0x55),
            "accent": RGBColor(0x1E, 0x3A, 0x8A),
            "stat_val": RGBColor(0x25, 0x63, 0xEB),
            "border": RGBColor(0xE2, 0xE8, 0xF0),
        }
    # Default dark theme
    return {
        "bg": RGBColor(0x0B, 0x0C, 0x10),
        "card_bg": RGBColor(0x14, 0x15, 0x17),
        "title": RGBColor(0xF1, 0xF5, 0xF9),
        "body": RGBColor(0x94, 0xA3, 0xB8),
        "accent": RGBColor(0x5B, 0x8D, 0xEF),
        "stat_val": RGBColor(0x60, 0xA5, 0xFA),
        "border": RGBColor(0x2E, 0x30, 0x36),
    }


def _create_slide_background(
    slide: Any, prs: Presentation, colors: dict[str, RGBColor]
):
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors["bg"]
    bg_shape.line.fill.background()
    return bg_shape


def _parse_markdown_to_slides(md_text: str) -> list[dict[str, Any]]:
    """Parse a Markdown text presentation into structured slide objects."""
    raw_slides = re.split(r"\n---\n|\n===\n", md_text)
    slides: list[dict[str, Any]] = []

    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        lines = [l.strip() for l in raw.split("\n") if l.strip()]
        if not lines:
            continue

        title = "Key Insights"
        content_lines: list[str] = []

        for line in lines:
            if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
                title = re.sub(r"^#+\s*", "", line).strip()
            else:
                content_lines.append(line)

        # Check for stats layout (lines like `> 99.9% : Uptime` or `**$50B** - Market Size` or `99.9% | Uptime`)
        stat_matches = []
        for cl in content_lines:
            stat_m = re.search(r"(\$?\d+[\d,\.]*[%kKmMbBtT]?|\b[A-Z0-9\.\+-]+%)\s*[:\|\-–—]\s*([A-Za-z0-9\s,\./]+)", cl)
            if stat_m:
                val = stat_m.group(1).strip()
                lbl = stat_m.group(2).strip()
                stat_matches.append({"value": val, "label": lbl})

        if len(stat_matches) >= 2 and len(stat_matches) == len(content_lines):
            slides.append({
                "type": "stats",
                "title": title,
                "stats": stat_matches[:4],
            })
            continue

        # Check for two-column layout (presence of two subheadings or distinct blocks)
        subheadings = [i for i, l in enumerate(content_lines) if l.startswith("#### ") or l.startswith("**") and l.endswith("**")]
        if len(subheadings) == 2:
            idx1, idx2 = subheadings[0], subheadings[1]
            head1 = content_lines[idx1].replace("**", "").replace("####", "").strip()
            head2 = content_lines[idx2].replace("**", "").replace("####", "").strip()
            b1 = [re.sub(r"^[-*\d\.]+\s*", "", l) for l in content_lines[idx1 + 1 : idx2] if l.startswith("- ") or l.startswith("* ") or re.match(r"^\d+\.", l)]
            b2 = [re.sub(r"^[-*\d\.]+\s*", "", l) for l in content_lines[idx2 + 1 :] if l.startswith("- ") or l.startswith("* ") or re.match(r"^\d+\.", l)]
            if b1 and b2:
                slides.append({
                    "type": "two_column",
                    "title": title,
                    "col1_title": head1,
                    "col1_bullets": b1,
                    "col2_title": head2,
                    "col2_bullets": b2,
                })
                continue

        # Default bullets layout
        bullets = []
        for cl in content_lines:
            clean_bullet = re.sub(r"^[-*>\d\.]+\s*", "", cl).strip()
            if clean_bullet:
                bullets.append(clean_bullet)

        slides.append({
            "type": "bullets",
            "title": title,
            "bullets": bullets or ["Key strategic takeaway."],
        })

    return slides


@tool
def create_powerpoint_presentation(
    filename: str,
    title: str = "Executive Briefing",
    subtitle: str = "",
    slides_json: str | None = None,
    content: str | None = None,
    theme: str = "dark",
) -> str:
    """Create a modern 16:9 widescreen PowerPoint (.pptx) presentation with formatted slide layouts.

    Supports either structured slides_json OR direct Markdown content string (slides delimited by ---).

    filename: Name of the file, e.g. "Platform_Pitch.pptx"
    title: Main title of the presentation (Title Slide)
    subtitle: Subtitle for the title slide
    theme: "dark" or "light" (default: "dark")
    content: Raw Markdown slides text with slides separated by '---'
    slides_json: Optional JSON string representing slides. Example:
    [
      {
        "type": "bullets",
        "title": "Platform Highlights",
        "bullets": ["Low latency", "Code Sandbox", "Deep integrations"]
      },
      {
        "type": "stats",
        "title": "Key Metrics",
        "stats": [{"value": "99.95%", "label": "Uptime"}, {"value": "<180ms", "label": "Latency"}]
      },
      {
        "type": "two_column",
        "title": "Comparison",
        "col1_title": "Before",
        "col1_bullets": ["Manual", "Slow"],
        "col2_title": "After",
        "col2_bullets": ["Automated", "Fast"]
      }
    ]
    """
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    slides_data: list[dict[str, Any]] = []

    # 1. Try parsing direct Markdown content if provided
    if content and isinstance(content, str) and content.strip():
        slides_data = _parse_markdown_to_slides(content)
    # 2. Try parsing slides_json
    elif slides_json and isinstance(slides_json, str) and slides_json.strip():
        try:
            parsed = json.loads(slides_json)
            if isinstance(parsed, list):
                slides_data = parsed
            elif isinstance(parsed, dict):
                slides_data = [parsed]
        except Exception:
            # Fallback: treat slides_json as raw markdown
            slides_data = _parse_markdown_to_slides(slides_json)

    if not slides_data:
        slides_data = [
            {
                "type": "bullets",
                "title": title,
                "bullets": ["Executive summary and strategic findings."],
            }
        ]

    prs = Presentation()
    # Set 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]
    colors = _get_theme_colors(theme)

    # 1. Title Slide
    title_slide = prs.slides.add_slide(blank_slide_layout)
    _create_slide_background(title_slide, prs, colors)

    title_box = title_slide.shapes.add_textbox(
        Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Calibri"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = colors["title"]
    p1.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Calibri"
        p2.font.size = Pt(20)
        p2.font.color.rgb = colors["accent"]
        p2.space_before = Pt(14)
        p2.alignment = PP_ALIGN.LEFT

    total_slides = 1

    # 2. Content Slides
    for slide_idx, s in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_slide_layout)
        _create_slide_background(slide, prs, colors)
        total_slides += 1

        stitle = s.get("title", f"Slide {slide_idx + 2}")
        stype = s.get("type", "bullets")

        # Slide Header Title
        header_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0)
        )
        htf = header_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = stitle
        hp.font.name = "Calibri"
        hp.font.size = Pt(28)
        hp.font.bold = True
        hp.font.color.rgb = colors["title"]

        # Slide Content Layouts
        if stype == "bullets":
            bullets = s.get("bullets", [])
            card_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.0),
                Inches(2.0),
                Inches(11.333),
                Inches(4.6),
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = colors["card_bg"]
            card_shape.line.color.rgb = colors["border"]

            ctf = card_shape.text_frame
            ctf.word_wrap = True
            ctf.margin_left = Inches(0.4)
            ctf.margin_top = Inches(0.4)

            for i, b in enumerate(bullets):
                bp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                bp.text = f"•   {b}"
                bp.font.name = "Calibri"
                bp.font.size = Pt(17)
                bp.font.color.rgb = colors["body"]
                bp.space_after = Pt(16)
                bp.line_spacing = 1.25

        elif stype == "stats":
            stats = s.get("stats", [])
            stat_count = min(len(stats), 4) or 1
            card_width = (11.333 - (stat_count - 1) * 0.4) / stat_count

            for i, stat in enumerate(stats[:4]):
                left = Inches(1.0 + i * (card_width + 0.4))
                stat_card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left,
                    Inches(2.2),
                    Inches(card_width),
                    Inches(4.0),
                )
                stat_card.fill.solid()
                stat_card.fill.fore_color.rgb = colors["card_bg"]
                stat_card.line.color.rgb = colors["border"]

                stf = stat_card.text_frame
                stf.word_wrap = True
                stf.margin_top = Inches(0.8)

                val_p = stf.paragraphs[0]
                val_p.text = str(stat.get("value", ""))
                val_p.font.name = "Calibri"
                val_p.font.size = Pt(36)
                val_p.font.bold = True
                val_p.font.color.rgb = colors["stat_val"]
                val_p.alignment = PP_ALIGN.CENTER

                lbl_p = stf.add_paragraph()
                lbl_p.text = str(stat.get("label", ""))
                lbl_p.font.name = "Calibri"
                lbl_p.font.size = Pt(14)
                lbl_p.font.color.rgb = colors["body"]
                lbl_p.alignment = PP_ALIGN.CENTER
                lbl_p.space_before = Pt(12)

        elif stype == "two_column":
            col_width = Inches(5.45)

            # Left Column Card
            col1_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.0),
                Inches(2.0),
                col_width,
                Inches(4.6),
            )
            col1_card.fill.solid()
            col1_card.fill.fore_color.rgb = colors["card_bg"]
            col1_card.line.color.rgb = colors["border"]

            tf1 = col1_card.text_frame
            tf1.word_wrap = True
            tf1.margin_left = Inches(0.35)
            tf1.margin_top = Inches(0.35)

            p_head1 = tf1.paragraphs[0]
            p_head1.text = s.get("col1_title", "Category A")
            p_head1.font.name = "Calibri"
            p_head1.font.size = Pt(18)
            p_head1.font.bold = True
            p_head1.font.color.rgb = colors["title"]
            p_head1.space_after = Pt(12)

            for b in s.get("col1_bullets", []):
                bp = tf1.add_paragraph()
                bp.text = f"•  {b}"
                bp.font.name = "Calibri"
                bp.font.size = Pt(14)
                bp.font.color.rgb = colors["body"]
                bp.space_after = Pt(8)

            # Right Column Card
            col2_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.85),
                Inches(2.0),
                col_width,
                Inches(4.6),
            )
            col2_card.fill.solid()
            col2_card.fill.fore_color.rgb = colors["card_bg"]
            col2_card.line.color.rgb = colors["accent"]

            tf2 = col2_card.text_frame
            tf2.word_wrap = True
            tf2.margin_left = Inches(0.35)
            tf2.margin_top = Inches(0.35)

            p_head2 = tf2.paragraphs[0]
            p_head2.text = s.get("col2_title", "Category B")
            p_head2.font.name = "Calibri"
            p_head2.font.size = Pt(18)
            p_head2.font.bold = True
            p_head2.font.color.rgb = colors["accent"]
            p_head2.space_after = Pt(12)

            for b in s.get("col2_bullets", []):
                bp = tf2.add_paragraph()
                bp.text = f"•  {b}"
                bp.font.name = "Calibri"
                bp.font.size = Pt(14)
                bp.font.color.rgb = colors["body"]
                bp.space_after = Pt(8)

    buffer = io.BytesIO()
    prs.save(buffer)
    pptx_bytes = buffer.getvalue()
    size_bytes = len(pptx_bytes)
    content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    from .s3_storage import upload_deliverable_to_s3

    s3_uri, download_url, fallback_data_uri = upload_deliverable_to_s3(
        file_bytes=pptx_bytes,
        filename=filename,
        content_type=content_type,
    )

    doc_event = {
        "filename": filename,
        "file_type": "powerpoint",
        "size_bytes": size_bytes,
        "s3_uri": s3_uri,
        "url": download_url,
        **({"data_uri": fallback_data_uri} if fallback_data_uri else {}),
        "summary": f"{title} ({total_slides} slides)",
        "title": title,
        "subtitle": subtitle,
        "theme": theme,
        "slides": slides_data,
    }

    from office_tools import document_queue

    document_queue.put_nowait(doc_event)

    return json.dumps(
        {
            "status": "success",
            "filename": filename,
            "file_type": "powerpoint",
            "size_bytes": size_bytes,
            "slides_count": total_slides,
            "s3_uri": s3_uri,
            "url": download_url,
            "download_ready": True,
            "summary": f"Successfully created PowerPoint presentation '{filename}' with {total_slides} slide(s).",
        },
        indent=2,
    )
